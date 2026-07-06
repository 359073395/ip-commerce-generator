from __future__ import annotations

import csv
import json
import re
import zipfile
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


BASE = Path(r"C:\Users\w'k'r\Documents\New project")
INVENTORY = BASE / "知识库" / "新学习资料" / "资产清单.json"
OUT_DIR = BASE / "知识库" / "新学习资料" / "素材" / "文档转写"
SKIP_PARTS = ["100套面试简历模板+大礼包"]
SUPPORTED = {".docx", ".xlsx", ".pptx", ".pdf", ".csv"}


def safe_name(text: str) -> str:
    text = re.sub(r"[\\/:*?\"<>|]+", "_", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text[:140] or "untitled"


def convert_docx(path: Path) -> str:
    try:
        doc = Document(str(path))
    except Exception:
        return convert_docx_xml_only(path)
    lines = [f"# {path.stem}", ""]
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
            lines.append("")
    for table_index, table in enumerate(doc.tables, start=1):
        lines.append(f"## 表格 {table_index}")
        rows = [[cell.text.strip().replace("\n", "<br>") for cell in row.cells] for row in table.rows]
        if rows:
            width = max(len(row) for row in rows)
            rows = [row + [""] * (width - len(row)) for row in rows]
            lines.append("| " + " | ".join(rows[0]) + " |")
            lines.append("| " + " | ".join(["---"] * width) + " |")
            for row in rows[1:]:
                lines.append("| " + " | ".join(row) + " |")
            lines.append("")
    return "\n".join(lines).strip() + "\n"


def convert_docx_xml_only(path: Path) -> str:
    lines = [f"# {path.stem}", "", "_注：该 DOCX 的媒体文件可能损坏，已使用 XML 正文兜底提取。_", ""]
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(path) as archive:
        xml = archive.read("word/document.xml")
    root = ET.fromstring(xml)
    for paragraph in root.findall(".//w:p", namespace):
        texts = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
        text = "".join(texts).strip()
        if text:
            lines.append(text)
            lines.append("")
    return "\n".join(lines).strip() + "\n"


def convert_xlsx(path: Path) -> str:
    wb = load_workbook(str(path), data_only=False, read_only=True)
    lines = [f"# {path.stem}", ""]
    for ws in wb.worksheets:
        lines.append(f"## Sheet: {ws.title}")
        rows = []
        for row in ws.iter_rows(values_only=True):
            values = ["" if value is None else str(value).replace("\n", "<br>") for value in row]
            if any(value.strip() for value in values):
                rows.append(values)
        if not rows:
            lines.append("")
            continue
        width = max(len(row) for row in rows)
        rows = [row + [""] * (width - len(row)) for row in rows]
        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("| " + " | ".join(["---"] * width) + " |")
        for row in rows[1:]:
            lines.append("| " + " | ".join(row) + " |")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def convert_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    lines = [f"# {path.stem}", ""]
    for slide_index, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {slide_index}")
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        if texts:
            for text in texts:
                lines.append(text)
                lines.append("")
        else:
            lines.append("_无可提取文本_")
            lines.append("")
    return "\n".join(lines).strip() + "\n"


def convert_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    lines = [f"# {path.stem}", ""]
    for page_index, page in enumerate(reader.pages, start=1):
        lines.append(f"## Page {page_index}")
        text = page.extract_text() or ""
        text = text.strip()
        lines.append(text if text else "_无可提取文本，可能是图片型PDF_")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def convert_csv(path: Path) -> str:
    lines = [f"# {path.stem}", ""]
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = list(csv.reader(handle))
    if not reader:
        return "\n".join(lines).strip() + "\n"
    width = max(len(row) for row in reader)
    rows = [row + [""] * (width - len(row)) for row in reader]
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * width) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines).strip() + "\n"


def should_skip(item: dict) -> bool:
    full_path = item["full_path"]
    return any(part in full_path for part in SKIP_PARTS)


def convert(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".docx":
        return convert_docx(path)
    if ext == ".xlsx":
        return convert_xlsx(path)
    if ext == ".pptx":
        return convert_pptx(path)
    if ext == ".pdf":
        return convert_pdf(path)
    if ext == ".csv":
        return convert_csv(path)
    raise ValueError(f"unsupported: {ext}")


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    items = json.loads(INVENTORY.read_text(encoding="utf-8"))
    report = []
    converted = 0
    skipped = 0
    failed = 0

    for item in items:
        ext = item["extension"].lower()
        if item["media_type"] != "document" or ext not in SUPPORTED or should_skip(item):
            skipped += 1
            continue
        src = Path(item["full_path"])
        out_name = safe_name(f"{item['id']}_{src.stem}") + ".md"
        out_path = OUT_DIR / out_name
        try:
            text = convert(src)
            header = [
                "---",
                f"id: {item['id']}",
                f"source: {src}",
                f"extension: {ext}",
                f"initial_category: {item['initial_category']}",
                "---",
                "",
            ]
            out_path.write_text("\n".join(header) + text, encoding="utf-8")
            report.append({"id": item["id"], "status": "converted", "source": str(src), "output": str(out_path)})
            converted += 1
        except Exception as exc:
            report.append({"id": item["id"], "status": "failed", "source": str(src), "error": str(exc)})
            failed += 1

    report_path = OUT_DIR / "_conversion_report.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"converted={converted} skipped={skipped} failed={failed}")
    print(report_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
